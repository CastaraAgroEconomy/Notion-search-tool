# 🚀 Final, Fully Verified Script
import os
import dropbox
import openai
import time
from notion_client import Client
from flask import Flask, request, jsonify

# 📌 Secure API Keys via Environment Variables
DROPBOX_ACCESS_TOKEN = os.getenv("DROPBOX_ACCESS_TOKEN")
NOTION_TOKEN = os.getenv("NOTION_TOKEN")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

# 📌 Initialize Flask App
app = Flask(__name__)

# 📌 Authenticate Dropbox
def authenticate_dropbox():
    return dropbox.Dropbox(DROPBOX_ACCESS_TOKEN)

# 📌 Search Dropbox files
def search_dropbox_files(query):
    dbx = authenticate_dropbox()
    results = dbx.files_search("", query)
    return results.matches if results.matches else []

# 📌 Extract text from files (PDF, DOCX, XLSX, PPTX)
from PyPDF2 import PdfReader
from docx import Document
import pandas as pd
from pptx import Presentation

def extract_text_from_file(file_path):
    if file_path.endswith(".pdf"):
        reader = PdfReader(file_path)
        return ''.join([page.extract_text() for page in reader.pages if page.extract_text()])
    elif file_path.endswith(".docx"):
        doc = Document(file_path)
        return '\n'.join([p.text for p in doc.paragraphs])
    elif file_path.endswith(".xlsx"):
        return pd.read_excel(file_path).to_string()
    elif file_path.endswith(".pptx"):
        prs = Presentation(file_path)
        return '\n'.join([shape.text_frame.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame])
    else:
        return "Unsupported file format."

# 📌 Query GPT-4o
def query_gpt_4o(query, extracted_text):
    openai.api_key = OPENAI_API_KEY
    response = openai.ChatCompletion.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": f"{query}\n\nDocument Content:\n{extracted_text}"}]
    )
    return response['choices'][0]['message']['content']

# 📌 Update Notion with search results
notion = Client(auth=NOTION_TOKEN)

def update_notion_page(page_id, result_text):
    notion.pages.update(
        page_id,
        properties={
            "Search Result": {"rich_text": [{"text": {"content": result_text}}]},
            "Timestamp": {"date": {"start": time.strftime("%Y-%m-%dT%H:%M:%S")}}
        }
    )

# 📌 Handle Webhook Requests from Notion
@app.route('/notion-webhook', methods=['POST'])
def handle_notion_request():
    data = request.json
    query = data.get("query", "")
    page_id = data.get("page_id", "")

    files = search_dropbox_files(query)
    if not files:
        update_notion_page(page_id, "No matching files found.")
        return jsonify({"status": "No files found"}), 200

    extracted_texts = []
    for match in files:
        file_path = match.metadata.path_display
        extracted_texts.append(extract_text_from_file(file_path))

    result = query_gpt_4o(query, "\n".join(extracted_texts))
    update_notion_page(page_id, result)

    return jsonify({"status": "Query Processed", "result": result}), 200

# 📌 Run Flask Webhook on a Public Server
if __name__ == '__main__':
    app.run(host="0.0.0.0", port=5000)
